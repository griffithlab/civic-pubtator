#!/usr/bin/env python3
import argparse, datetime, json, os, shutil, subprocess, sys, tempfile

# ── LibreOffice detection ─────────────────────────────────────────────────────

def find_soffice():
    path = shutil.which("soffice")
    if path:
        return path
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.isfile(mac_path):
        return mac_path
    return None

def soffice_convert(soffice, src, out_dir, convert_to="pdf"):
    """Convert src to convert_to format in out_dir via soffice.

    Uses --norestore and a per-invocation --user-installation temp directory
    so each call runs as a fully independent process.  Without this, LibreOffice's
    single-instance check causes a second call to hand off to the still-running
    instance from a previous conversion and exit immediately with code 1.

    Returns the path of the created output file.
    """
    # Set LD_LIBRARY_PATH so soffice.bin can find its internal libraries on Linux
    # without requiring a system-wide ldconfig entry (which causes conflicts).
    env = os.environ.copy()
    env['LD_LIBRARY_PATH'] = '/usr/lib/libreoffice/program:' + env.get('LD_LIBRARY_PATH', '')
    user_install = tempfile.mkdtemp(prefix="soffice-user-")
    try:
        result = subprocess.run(
            [
                soffice, "--headless", "--norestore",
                f"-env:UserInstallation=file://{user_install}",
                "--convert-to", convert_to, "--outdir", out_dir, src,
            ],
            capture_output=True, text=True, env=env,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"soffice failed (exit {result.returncode}):\n"
                f"STDERR: {result.stderr.strip()}\nSTDOUT: {result.stdout.strip()}"
            )
    finally:
        shutil.rmtree(user_install, ignore_errors=True)
    stem = os.path.splitext(os.path.basename(src))[0]
    ext = convert_to.split(":")[0]  # e.g. "pdf" or "xlsx" (strip filter hints if any)
    return os.path.join(out_dir, f"{stem}.{ext}")

# ── Conversion sidecar ───────────────────────────────────────────────────────

def _write_conversion_sidecar(out_dir, stem, source_path, method):
    """Write a .conversion.json sidecar alongside the converted PDF.

    Records the original source filename and the method used so that
    civic_pubtator.py can include this information in content_capture_stats.tsv.
    """
    sidecar = {
        "source_file":       os.path.basename(source_path),
        "conversion_method": method,
    }
    with open(os.path.join(out_dir, stem + ".conversion.json"), "w") as fh:
        json.dump(sidecar, fh)


# ── Per-type processors ───────────────────────────────────────────────────────

def process_pdf(src, stem, s_dir):
    out_dir = os.path.join(s_dir, stem)
    os.makedirs(out_dir, exist_ok=True)
    dst = os.path.join(out_dir, stem + ".pdf")
    shutil.copy2(src, dst)
    print(f"  Copied → {dst}")
    _write_conversion_sidecar(out_dir, stem, src, "direct")

def process_word(src, stem, s_dir, soffice):
    out_dir = os.path.join(s_dir, stem)
    os.makedirs(out_dir, exist_ok=True)

    ext = os.path.splitext(src)[1].lower()
    tmp_dir = None

    try:
        # .doc → .docx first so the python-docx fallback is always reachable
        if ext == ".doc":
            if not soffice:
                raise RuntimeError(".doc requires LibreOffice; install it and retry")
            tmp_dir = tempfile.mkdtemp()
            src = soffice_convert(soffice, src, tmp_dir, convert_to="docx")
            print(f"  .doc → .docx (soffice): {src}")

        if soffice:
            try:
                created = soffice_convert(soffice, src, out_dir)
                dst = os.path.join(out_dir, stem + ".pdf")
                if created != dst:
                    os.replace(created, dst)
                print(f"  Word → PDF (soffice): {dst}")
                _write_conversion_sidecar(out_dir, stem, src, "soffice")
                return
            except RuntimeError as exc:
                print(f"  WARNING: soffice failed — {exc}")
                print(f"  Falling back to python-docx + reportlab")

        # Fallback: python-docx + reportlab (.doc was already converted to .docx above)
        try:
            from docx import Document
        except ImportError:
            sys.exit("ERROR: python-docx not installed. Run: pip3 install python-docx")
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
        except ImportError:
            sys.exit("ERROR: reportlab not installed. Run: pip3 install reportlab")

        dst = os.path.join(out_dir, stem + ".pdf")
        doc = Document(src)
        styles = getSampleStyleSheet()
        story = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                story.append(Paragraph(text, styles["Normal"]))
                story.append(Spacer(1, 4))
        SimpleDocTemplate(dst, pagesize=letter).build(story)
        print(f"  Word → PDF (python-docx fallback): {dst}")
        _write_conversion_sidecar(out_dir, stem, src, "python-docx")
    finally:
        if tmp_dir:
            shutil.rmtree(tmp_dir, ignore_errors=True)

def process_powerpoint(src, stem, s_dir, soffice):
    out_dir = os.path.join(s_dir, stem)
    os.makedirs(out_dir, exist_ok=True)

    ext = os.path.splitext(src)[1].lower()
    tmp_dir = None

    try:
        # .ppt → .pptx first so the python-pptx fallback is always reachable
        if ext == ".ppt":
            if not soffice:
                raise RuntimeError(".ppt requires LibreOffice; install it and retry")
            tmp_dir = tempfile.mkdtemp()
            src = soffice_convert(soffice, src, tmp_dir, convert_to="pptx")
            print(f"  .ppt → .pptx (soffice): {src}")

        if soffice:
            try:
                created = soffice_convert(soffice, src, out_dir)
                dst = os.path.join(out_dir, stem + ".pdf")
                if created != dst:
                    os.replace(created, dst)
                print(f"  PowerPoint → PDF (soffice): {dst}")
                _write_conversion_sidecar(out_dir, stem, src, "soffice")
                return
            except RuntimeError as exc:
                print(f"  WARNING: soffice failed — {exc}")
                print(f"  Falling back to python-pptx + reportlab")

        # Fallback: python-pptx + reportlab (.ppt was already converted to .pptx above)
        try:
            from pptx import Presentation
        except ImportError:
            sys.exit("ERROR: python-pptx not installed. Run: pip3 install python-pptx")
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
        except ImportError:
            sys.exit("ERROR: reportlab not installed. Run: pip3 install reportlab")

        dst = os.path.join(out_dir, stem + ".pdf")
        prs = Presentation(src)
        styles = getSampleStyleSheet()
        story = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            story.append(Paragraph(f"<b>Slide {slide_num}</b>", styles["Heading2"]))
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            story.append(Paragraph(text, styles["Normal"]))
                            story.append(Spacer(1, 4))
            story.append(Spacer(1, 12))
        SimpleDocTemplate(dst, pagesize=letter).build(story)
        print(f"  PowerPoint → PDF (python-pptx fallback): {dst}")
        _write_conversion_sidecar(out_dir, stem, src, "python-pptx")
    finally:
        if tmp_dir:
            shutil.rmtree(tmp_dir, ignore_errors=True)


def _numeric_fraction(cells):
    """Return fraction of non-empty cells that parse as a number."""
    non_empty = [v for v in cells if v is not None and str(v).strip() != ""]
    if not non_empty:
        return 0.0
    count = 0
    for v in non_empty:
        if isinstance(v, (int, float)) and not isinstance(v, bool):
            count += 1
        else:
            try:
                float(str(v).strip().rstrip("%"))
                count += 1
            except (ValueError, TypeError):
                pass
    return count / len(non_empty)


def _filter_numeric_columns(rows, threshold=0.95):
    """
    Drop columns whose data rows (excluding the header) are >= threshold numeric.

    Returns (filtered_rows, dropped) where dropped is a list of
    (original_col_index, col_name, numeric_fraction) for each removed column.
    """
    if len(rows) < 2:
        return rows, []

    header = rows[0]
    data_rows = rows[1:]
    ncols = max(len(r) for r in rows)

    keep, dropped = [], []
    for col_idx in range(ncols):
        cells = [r[col_idx] if col_idx < len(r) else None for r in data_rows]
        frac = _numeric_fraction(cells)
        if frac >= threshold:
            col_name = (
                header[col_idx]
                if col_idx < len(header) and header[col_idx] is not None
                else f"col_{col_idx + 1}"
            )
            dropped.append((col_idx, str(col_name), frac))
        else:
            keep.append(col_idx)

    if not dropped:
        return rows, []

    filtered = [
        tuple(r[i] if i < len(r) else None for i in keep)
        for r in rows
    ]
    return filtered, dropped


def process_excel(src, stem, s_dir, soffice, max_rows, max_tabs=15):
    original_src = src  # preserve before any .xls → .xlsx reassignment
    ext = os.path.splitext(src)[1].lower()
    tmp_xls_dir = None
    if not max_rows:
        max_rows = float('inf')
    if not max_tabs:
        max_tabs = float('inf')

    # use_soffice tracks whether soffice is available AND has not crashed yet for this file.
    use_soffice = bool(soffice)

    if ext == ".xls":
        if use_soffice:
            tmp_xls_dir = tempfile.mkdtemp()
            try:
                src = soffice_convert(soffice, src, tmp_xls_dir, convert_to="xlsx")
                print(f"  .xls → .xlsx (soffice): {src}")
            except RuntimeError as exc:
                shutil.rmtree(tmp_xls_dir, ignore_errors=True)
                tmp_xls_dir = None
                print(f"  WARNING: soffice failed for .xls→.xlsx — {exc}")
                print(f"  Falling back to xlrd")
                use_soffice = False

        if not use_soffice:
            # fallback: xlrd reads .xls; write to temp .xlsx so openpyxl can load it below
            try:
                import xlrd
            except ImportError:
                sys.exit("ERROR: .xls requires LibreOffice or xlrd. Run: pip3 install xlrd")
            try:
                import openpyxl as _openpyxl
            except ImportError:
                sys.exit("ERROR: openpyxl not installed. Run: pip3 install openpyxl")
            tmp_xls_dir = tempfile.mkdtemp()
            wb_xls = xlrd.open_workbook(src)
            wb_out = _openpyxl.Workbook()
            wb_out.remove(wb_out.active)
            for sname in wb_xls.sheet_names():
                ws_in = wb_xls.sheet_by_name(sname)
                ws_out = wb_out.create_sheet(title=sname)
                for i in range(min(ws_in.nrows, max_rows)):
                    ws_out.append(ws_in.row_values(i))
            wb_xls.release_resources()
            tmp_xlsx = os.path.join(tmp_xls_dir, stem + ".xlsx")
            wb_out.save(tmp_xlsx)
            wb_out.close()
            src = tmp_xlsx
            print(f"  .xls → .xlsx (xlrd): {src}")

    try:
        import openpyxl
    except ImportError:
        sys.exit("ERROR: openpyxl not installed. Run: pip3 install openpyxl")

    wb = openpyxl.load_workbook(src, data_only=True)

    for tab_num, sheet_name in enumerate(wb.sheetnames, start=1):
        if tab_num > max_tabs:
            remaining = len(wb.sheetnames) - int(max_tabs)
            print(f"  Tab limit ({int(max_tabs)}) reached: skipping {remaining} remaining tab(s)")
            break
        ws = wb[sheet_name]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            rows.append(row)

        if not rows:
            print(f"  Sheet '{sheet_name}' (tab {tab_num}): empty, skipping")
            continue

        original_ncols = max(len(r) for r in rows)
        rows, dropped_cols = _filter_numeric_columns(rows)
        for col_idx, col_name, frac in dropped_cols:
            print(f"  Dropping column '{col_name}' (col {col_idx + 1}): {frac:.0%} numeric")
        if dropped_cols:
            kept = original_ncols - len(dropped_cols)
            print(f"  → {kept}/{original_ncols} columns kept after numeric filter")

        tab_label = f"tab_{tab_num:02d}"
        out_dir = os.path.join(s_dir, stem, tab_label)
        os.makedirs(out_dir, exist_ok=True)
        dst = os.path.join(out_dir, f"{stem}.pdf")

        soffice_tab_ok = False
        if use_soffice:
            from openpyxl.utils import get_column_letter
            tmp_dir = tempfile.mkdtemp()
            try:
                tmp_wb = openpyxl.Workbook()
                tmp_ws = tmp_wb.active
                tmp_ws.title = sheet_name
                for row in rows:
                    tmp_ws.append([v if v is not None else "" for v in row])

                # Auto-size columns based on content
                ncols = max(len(r) for r in rows)
                for col_idx in range(1, ncols + 1):
                    max_len = max(
                        len(str(r[col_idx - 1])) if col_idx <= len(r) and r[col_idx - 1] is not None else 0
                        for r in rows
                    )
                    tmp_ws.column_dimensions[get_column_letter(col_idx)].width = max(12, min(max_len + 2, 50))

                # Landscape, A3, fit all columns onto one page wide
                tmp_ws.page_setup.orientation = "landscape"
                tmp_ws.page_setup.paperSize = 8  # A3
                tmp_ws.page_setup.fitToWidth = 1
                tmp_ws.page_setup.fitToHeight = 0
                tmp_ws.sheet_properties.pageSetUpPr.fitToPage = True

                # Header: sheet name centred, tab number on the right
                tmp_ws.oddHeader.center.text = sheet_name
                tmp_ws.oddHeader.right.text  = f"Tab {tab_num}"

                tmp_xlsx = os.path.join(tmp_dir, f"{stem}.xlsx")
                tmp_wb.save(tmp_xlsx)
                created = soffice_convert(soffice, tmp_xlsx, tmp_dir)
                shutil.move(created, dst)
                soffice_tab_ok = True
                print(f"  Sheet '{sheet_name}' (tab {tab_num}, {len(rows)} rows) → {dst} (soffice)")
            except RuntimeError as exc:
                print(f"  WARNING: soffice failed for tab {tab_num} — {exc}")
                print(f"  Falling back to reportlab")
            finally:
                shutil.rmtree(tmp_dir, ignore_errors=True)

        if not use_soffice or not soffice_tab_ok:
            # ── Fallback: reportlab ───────────────────────────────────────────
            try:
                from reportlab.lib import colors
                from reportlab.lib.pagesizes import landscape, letter
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.lib.units import inch
                from reportlab.platypus import SimpleDocTemplate, Spacer, Table, TableStyle, Paragraph
            except ImportError:
                sys.exit("ERROR: reportlab not installed. Run: pip3 install reportlab")

            MARGINS = 0.5 * inch
            MIN_COL_WIDTH = 1.3 * inch

            str_rows = [[str(v) if v is not None else "" for v in row] for row in rows]
            ncols = max(len(r) for r in str_rows)
            str_rows = [r + [""] * (ncols - len(r)) for r in str_rows]

            col_width  = max(MIN_COL_WIDTH, landscape(letter)[0] / ncols)
            page_width = max(landscape(letter)[0], ncols * col_width + 2 * MARGINS)
            page_height = landscape(letter)[1]
            font_size  = max(5, min(8, int(col_width / 0.12)))

            styles = getSampleStyleSheet()
            title = Paragraph(f"<b>{sheet_name}</b>  <font size='8'>(Tab {tab_num})</font>",
                              styles["Normal"])

            table = Table(str_rows, colWidths=[col_width] * ncols, repeatRows=1)
            table.setStyle(TableStyle([
                ("BACKGROUND",     (0, 0), (-1,  0), colors.grey),
                ("TEXTCOLOR",      (0, 0), (-1,  0), colors.whitesmoke),
                ("FONTSIZE",       (0, 0), (-1, -1), font_size),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
                ("GRID",           (0, 0), (-1, -1), 0.25, colors.black),
                ("VALIGN",         (0, 0), (-1, -1), "TOP"),
                ("WORDWRAP",       (0, 0), (-1, -1), True),
            ]))
            SimpleDocTemplate(dst, pagesize=(page_width, page_height),
                              leftMargin=MARGINS, rightMargin=MARGINS,
                              topMargin=MARGINS, bottomMargin=MARGINS).build(
                [title, Spacer(1, 0.15 * inch), table])
            print(f"  Sheet '{sheet_name}' (tab {tab_num}, {len(rows)} rows) → {dst} (reportlab fallback)")

        _write_conversion_sidecar(out_dir, stem, original_src,
                                  "soffice" if soffice_tab_ok else "reportlab")

    wb.close()
    if tmp_xls_dir:
        shutil.rmtree(tmp_xls_dir, ignore_errors=True)

# ── Main ──────────────────────────────────────────────────────────────────────

IGNORED_FILES = {".DS_Store"}

def main():
    parser = argparse.ArgumentParser(
        description="Convert supplementary files in <input_dir>/s/ to PDFs."
    )
    parser.add_argument("input_dir",
                        help="Source directory containing an s/ subdirectory")
    parser.add_argument("--no-libreoffice", action="store_true",
                        help="Use the reportlab/python-docx fallback even if LibreOffice is installed")
    parser.add_argument("--max-rows", type=int, default=1000, metavar="N",
                        help="Maximum rows to read per Excel sheet tab when converting "
                             ".xls/.xlsx to PDF (default: 1000; use 0 for no limit). "
                             "This caps the size of each converted table. In the main "
                             "pipeline a second filter, --max-chars, removes any document "
                             "whose total converted text still exceeds the character limit "
                             "after this row cap is applied.")
    parser.add_argument("--max-tabs", type=int, default=15, metavar="N",
                        help="Maximum number of tabs to extract from a single Excel "
                             "spreadsheet (default: 15; use 0 for no limit). Tabs beyond "
                             "this limit are skipped with a warning.")
    args = parser.parse_args()

    if args.no_libreoffice:
        soffice = None
        print("LibreOffice disabled by --no-libreoffice; using reportlab/python-docx fallback.")
    else:
        soffice = find_soffice()
        if soffice:
            print(f"LibreOffice found: {soffice}")
        else:
            print("WARNING: LibreOffice (soffice) not found.")
            print("         Word and Excel conversion will use a basic fallback with reduced fidelity.")
            print("         For best results install LibreOffice:")
            print("           macOS:  brew install --cask libreoffice")
            print("           Ubuntu: sudo apt-get install -y libreoffice\n")

    input_dir = os.path.abspath(args.input_dir)
    s_dir = os.path.join(input_dir, "s")

    if not os.path.isdir(s_dir):
        print(f"No s/ subdirectory found in {input_dir}, nothing to do.")
        sys.exit(0)

    failure_log = os.path.join(s_dir, "CONVERSION_FAILURES.log")

    for fname in sorted(os.listdir(s_dir)):
        fpath = os.path.join(s_dir, fname)
        if not os.path.isfile(fpath) or fname in IGNORED_FILES or fname.startswith("~$"):
            continue

        stem = os.path.splitext(fname)[0]
        ext  = os.path.splitext(fname)[1].lower()

        print(f"\nProcessing: {fname}")

        try:
            if ext == ".pdf":
                process_pdf(fpath, stem, s_dir)
            elif ext in (".docx", ".doc"):
                process_word(fpath, stem, s_dir, soffice)
            elif ext in (".xlsx", ".xls"):
                process_excel(fpath, stem, s_dir, soffice, args.max_rows, args.max_tabs)
            elif ext in (".pptx", ".ppt"):
                process_powerpoint(fpath, stem, s_dir, soffice)
            else:
                print(f"  Unsupported type ({ext}), skipping.")
        except Exception as exc:
            msg = f"{exc}"
            print(f"  ERROR: could not convert {fname} (all methods failed) — {msg}")
            ts = datetime.datetime.now().isoformat(timespec="seconds")
            with open(failure_log, "a") as fh:
                fh.write(f"{ts}  {fname}: {msg}\n")
            print(f"  Recorded in {failure_log}")

if __name__ == "__main__":
    main()
